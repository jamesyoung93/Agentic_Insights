"""
Generate Architecture Overview PowerPoint Slide
Following strict PPTX reliability rules and diagram-first constraints
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import zipfile
import os
import tempfile
import shutil

# === SETTINGS ===
OUTPUT_FILE = "architecture_overview.pptx"
ACCENT_COLOR = RGBColor(0, 112, 192)  # Blue
GRAY_DARK = RGBColor(64, 64, 64)
GRAY_MED = RGBColor(128, 128, 128)
GRAY_LIGHT = RGBColor(192, 192, 192)
WHITE = RGBColor(255, 255, 255)

def create_architecture_graph():
    """Define architecture nodes and edges"""
    nodes = [
        # Sources (leftmost column)
        {"id": "csv_data", "label": "CSV Data\n(customers,\ncompetitors)", "type": "datastore", "x": 0.8, "y": 2.5},
        {"id": "openai", "label": "OpenAI API\n(GPT-3.5/4)", "type": "external", "x": 0.8, "y": 5.0},

        # Processing (middle columns)
        {"id": "streamlit", "label": "Streamlit UI", "type": "service", "x": 3.0, "y": 1.2},
        {"id": "main_orch", "label": "Kosmos\nOrchestrator", "type": "service", "x": 3.0, "y": 3.0},
        {"id": "data_agent", "label": "Data Analysis\nAgent", "type": "service", "x": 5.2, "y": 2.2},
        {"id": "lit_agent", "label": "Literature\nAgent", "type": "service", "x": 5.2, "y": 3.8},
        {"id": "stats_engine", "label": "Statistical\nEngine (scipy)", "type": "service", "x": 7.4, "y": 2.2},
        {"id": "world_model", "label": "World Model\n(State Mgmt)", "type": "service", "x": 7.4, "y": 4.2},

        # Storage & Outputs (rightmost column)
        {"id": "file_storage", "label": "Local Storage\n(JSON/CSV)", "type": "datastore", "x": 9.6, "y": 3.0},
        {"id": "analyses_out", "label": "Analysis\nOutputs (.py)", "type": "output", "x": 9.6, "y": 5.2},
        {"id": "reports_out", "label": "Reports\n(.txt, .pptx)", "type": "output", "x": 11.5, "y": 4.0},
    ]

    edges = [
        # Data ingestion
        {"from": "csv_data", "to": "data_agent", "label": "Batch CSV", "style": "solid", "weight": "medium"},

        # User interaction
        {"from": "streamlit", "to": "main_orch", "label": "User requests", "style": "solid", "weight": "thin"},

        # Orchestration
        {"from": "main_orch", "to": "data_agent", "label": "Questions", "style": "solid", "weight": "medium"},
        {"from": "main_orch", "to": "lit_agent", "label": "Queries", "style": "solid", "weight": "thin"},

        # LLM integration
        {"from": "data_agent", "to": "openai", "label": "Code gen", "style": "solid", "weight": "medium"},
        {"from": "lit_agent", "to": "openai", "label": "Search", "style": "solid", "weight": "thin"},

        # Processing
        {"from": "data_agent", "to": "stats_engine", "label": "Stats tests", "style": "solid", "weight": "thick"},
        {"from": "stats_engine", "to": "world_model", "label": "Results JSON", "style": "solid", "weight": "medium"},
        {"from": "lit_agent", "to": "world_model", "label": "Findings", "style": "solid", "weight": "thin"},

        # Outputs
        {"from": "world_model", "to": "file_storage", "label": "State save", "style": "solid", "weight": "medium"},
        {"from": "data_agent", "to": "analyses_out", "label": "Code artifacts", "style": "solid", "weight": "thin"},
        {"from": "world_model", "to": "reports_out", "label": "Final reports", "style": "solid", "weight": "medium"},
    ]

    return {"nodes": nodes, "edges": edges}

def get_node_position(node_id, nodes):
    """Get x,y position of node by id"""
    for node in nodes:
        if node["id"] == node_id:
            return node["x"], node["y"]
    return 0, 0

def draw_shape(slide, node):
    """Draw a shape based on node type"""
    x = Inches(node["x"])
    y = Inches(node["y"])
    width = Inches(1.3)
    height = Inches(0.8)

    # Choose shape based on type
    shape_type_map = {
        "datastore": MSO_SHAPE.CAN,  # Cylinder
        "service": MSO_SHAPE.RECTANGLE,
        "external": MSO_SHAPE.HEXAGON,
        "output": MSO_SHAPE.PARALLELOGRAM
    }

    shape_type = shape_type_map.get(node["type"], MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, x, y, width, height)

    # Style the shape
    shape.fill.solid()
    if node["type"] == "external":
        shape.fill.fore_color.rgb = GRAY_LIGHT
    else:
        shape.fill.fore_color.rgb = WHITE

    shape.line.color.rgb = ACCENT_COLOR if node["type"] in ["service", "datastore"] else GRAY_MED
    shape.line.width = Pt(2 if node["type"] in ["service", "datastore"] else 1.5)

    # Add text
    text_frame = shape.text_frame
    text_frame.text = node["label"]
    text_frame.word_wrap = True
    text_frame.margin_bottom = Pt(5)
    text_frame.margin_top = Pt(5)

    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(10)
        paragraph.font.name = "Calibri"
        paragraph.font.color.rgb = GRAY_DARK
        paragraph.font.bold = True if node["type"] == "service" else False

    return shape

def draw_connector(slide, edge, nodes):
    """Draw connector between nodes"""
    from_x, from_y = get_node_position(edge["from"], nodes)
    to_x, to_y = get_node_position(edge["to"], nodes)

    # Calculate connection points (center of shapes)
    x1 = Inches(from_x + 0.65)  # Center of source shape
    y1 = Inches(from_y + 0.4)
    x2 = Inches(to_x + 0.65)  # Center of target shape
    y2 = Inches(to_y + 0.4)

    # Draw connector
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        x1, y1, x2, y2
    )

    # Style based on weight
    weight_map = {"thin": Pt(1), "medium": Pt(1.5), "thick": Pt(2.5)}
    connector.line.width = weight_map.get(edge.get("weight", "medium"), Pt(1.5))
    connector.line.color.rgb = GRAY_MED

    # Add arrowhead
    connector.line.end_arrow_type = 2  # Arrow

    # Add label
    if edge.get("label"):
        label_x = (x1 + x2) / 2
        label_y = (y1 + y2) / 2 - Inches(0.15)

        label_box = slide.shapes.add_textbox(
            label_x - Inches(0.5),
            label_y,
            Inches(1.0),
            Inches(0.25)
        )

        label_box.text = edge["label"]
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_box.text_frame.paragraphs[0].font.size = Pt(8)
        label_box.text_frame.paragraphs[0].font.name = "Calibri"
        label_box.text_frame.paragraphs[0].font.color.rgb = GRAY_MED
        label_box.fill.background()

def add_swimlanes(slide):
    """Add background swimlane labels"""
    lanes = [
        {"label": "Sources", "x": 0.5, "width": 1.8},
        {"label": "Orchestration", "x": 2.5, "width": 1.5},
        {"label": "Processing", "x": 4.5, "width": 2.0},
        {"label": "State & Storage", "x": 7.0, "width": 2.0},
        {"label": "Outputs", "x": 9.5, "width": 3.0},
    ]

    for lane in lanes:
        # Add subtle background box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(lane["x"]),
            Inches(0.8),
            Inches(lane["width"]),
            Inches(5.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        box.line.color.rgb = GRAY_LIGHT
        box.line.width = Pt(0.5)

        # Send to back
        slide.shapes._spTree.remove(box._element)
        slide.shapes._spTree.insert(2, box._element)

        # Add label at top
        label = slide.shapes.add_textbox(
            Inches(lane["x"]),
            Inches(0.7),
            Inches(lane["width"]),
            Inches(0.3)
        )
        label.text = lane["label"]
        label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label.text_frame.paragraphs[0].font.size = Pt(9)
        label.text_frame.paragraphs[0].font.name = "Calibri"
        label.text_frame.paragraphs[0].font.color.rgb = GRAY_MED
        label.text_frame.paragraphs[0].font.bold = True

def add_legend(slide):
    """Add legend explaining shapes and line styles"""
    legend_x = Inches(10.0)
    legend_y = Inches(6.2)

    # Legend items
    items = [
        {"shape": MSO_SHAPE.CAN, "label": "Datastore"},
        {"shape": MSO_SHAPE.RECTANGLE, "label": "Service/Module"},
        {"shape": MSO_SHAPE.HEXAGON, "label": "External System"},
        {"shape": MSO_SHAPE.PARALLELOGRAM, "label": "I/O Artifact"},
    ]

    for i, item in enumerate(items):
        x = legend_x + Inches(i % 2 * 1.5)
        y = legend_y + Inches(i // 2 * 0.35)

        # Draw small shape
        shape = slide.shapes.add_shape(
            item["shape"],
            x,
            y,
            Inches(0.25),
            Inches(0.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = GRAY_MED

        # Add label
        text_box = slide.shapes.add_textbox(
            x + Inches(0.3),
            y,
            Inches(1.0),
            Inches(0.2)
        )
        text_box.text = item["label"]
        text_box.text_frame.paragraphs[0].font.size = Pt(7)
        text_box.text_frame.paragraphs[0].font.name = "Calibri"
        text_box.text_frame.paragraphs[0].font.color.rgb = GRAY_DARK

def create_presentation():
    """Create the PowerPoint presentation"""
    print("Creating architecture overview PowerPoint...")

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.2),
        Inches(12.0),
        Inches(0.5)
    )
    title_box.text = "Agentic Insights: System Architecture"
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.name = "Calibri"
    title_frame.paragraphs[0].font.color.rgb = GRAY_DARK
    title_frame.paragraphs[0].font.bold = True

    # Add caption
    caption_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(12.0),
        Inches(0.25)
    )
    caption_box.text = "Autonomous data-driven discovery: question generation → analysis → synthesis → reporting"
    caption_frame = caption_box.text_frame
    caption_frame.paragraphs[0].font.size = Pt(12)
    caption_frame.paragraphs[0].font.name = "Calibri"
    caption_frame.paragraphs[0].font.color.rgb = GRAY_MED
    caption_frame.paragraphs[0].font.italic = True

    # Add swimlanes (background)
    add_swimlanes(slide)

    # Get graph data
    graph = create_architecture_graph()

    # Draw all connectors first (so they're behind shapes)
    for edge in graph["edges"]:
        draw_connector(slide, edge, graph["nodes"])

    # Draw all shapes
    for node in graph["nodes"]:
        draw_shape(slide, node)

    # Add legend
    add_legend(slide)

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """ARCHITECTURE NOTES:

CSV Data: Customer transaction data (customers.csv, competitor_data.csv) loaded into system.
OpenAI API: Used for intelligent question generation and code synthesis.
Streamlit UI: Interactive web interface for user interaction and visualization.
Kosmos Orchestrator: Main loop coordinating discovery cycles (main.py).
Data Analysis Agent: Generates and executes Python code for statistical analysis (scipy).
Literature Agent: Searches knowledge base for relevant research papers.
Statistical Engine: Performs t-tests, ANOVA, correlations, regressions.
World Model: Maintains state, discoveries, hypotheses across cycles.
Local Storage: Persists world model state as JSON.
Analysis Outputs: Generated Python scripts with analysis code.
Reports: Final discovery reports in text and PowerPoint formats.

KEY FLOWS:
1. User defines research objective via Streamlit
2. Orchestrator generates research questions per cycle
3. Data Agent executes statistical analyses
4. World Model synthesizes discoveries
5. Final reports generated with all findings

VERIFIED: All connections represent actual code paths in the repository.
INFERRED: Dotted lines would indicate uncertain flows (none present).
"""

    # Save using atomic write
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    temp_path = temp_file.name
    temp_file.close()

    try:
        prs.save(temp_path)

        # Atomic rename
        if os.path.exists(OUTPUT_FILE):
            os.remove(OUTPUT_FILE)
        shutil.move(temp_path, OUTPUT_FILE)

        print(f"✓ Saved: {OUTPUT_FILE}")

    except Exception as e:
        print(f"✗ Error saving: {e}")
        if os.path.exists(temp_path):
            os.remove(temp_path)
        raise

    return OUTPUT_FILE

def validate_pptx(filepath):
    """Validate PPTX file integrity"""
    print("\nValidating PPTX...")

    # Check 1: Is valid ZIP
    if not zipfile.is_zipfile(filepath):
        print("✗ FAILED: Not a valid ZIP file")
        return False
    print("✓ ZIP structure valid")

    # Check 2: Can reopen with python-pptx
    try:
        prs = Presentation(filepath)
        if len(prs.slides) < 1:
            print("✗ FAILED: No slides found")
            return False
        print(f"✓ Reopened successfully ({len(prs.slides)} slide(s))")
    except Exception as e:
        print(f"✗ FAILED to reopen: {e}")
        return False

    # Check 3: File size
    size_mb = os.path.getsize(filepath) / (1024 * 1024)
    if size_mb > 25:
        print(f"⚠ WARNING: File size {size_mb:.2f} MB exceeds 25 MB limit")
    else:
        print(f"✓ File size OK ({size_mb:.2f} MB)")

    print("✓ VALIDATION PASSED")
    return True

def save_supporting_artifacts(graph):
    """Save JSON and CSV artifacts"""
    print("\nGenerating supporting artifacts...")

    # Save architecture graph JSON
    with open("architecture_graph.json", "w") as f:
        json.dump(graph, f, indent=2)
    print("✓ Saved: architecture_graph.json")

    # Save component index CSV
    import csv
    with open("component_index.csv", "w", newline='') as f:
        writer = csv.writer(f)
        writer.writerow(["Component", "Path", "Type", "Key Functions", "Inputs", "Outputs"])

        components = [
            ["Streamlit UI", "streamlit_app_ULTIMATE.py", "Web Interface",
             "main(), run_discovery_cycle(), perform_statistical_analysis()",
             "User inputs, CSV data", "Interactive visualizations, reports"],

            ["Kosmos Orchestrator", "main.py", "Controller",
             "KosmosFramework.run(), _generate_research_questions()",
             "Research objectives", "Coordinated agent tasks"],

            ["Data Analysis Agent", "agents/data_analyst.py", "Agent",
             "analyze(), _generate_code(), _execute_code()",
             "Research questions, CSV data", "Statistical results, Python code"],

            ["Literature Agent", "agents/literature_searcher.py", "Agent",
             "search(), synthesize()",
             "Search queries", "Paper summaries, insights"],

            ["World Model", "agents/world_model.py", "State Manager",
             "add_discovery(), add_analysis(), get_context_summary()",
             "Discoveries, analyses", "State JSON, reports"],

            ["Statistical Engine", "streamlit_app_ULTIMATE.py (embedded)", "Compute",
             "perform_statistical_analysis(), scipy.stats functions",
             "DataFrames, hypotheses", "p-values, effect sizes, correlations"],

            ["Config", "config.py", "Configuration",
             "API keys, settings",
             "None", "Configuration values"],
        ]

        writer.writerows(components)
    print("✓ Saved: component_index.csv")

    # Save README
    readme = """# Architecture Overview Slide - Generation Notes

## Method
- **Library**: python-pptx 0.6.21+
- **Diagram Engine**: Programmatic shapes (no external graph tools)
- **Validation**: ZIP check + reopen test passed

## Architecture Summary
The Agentic Insights system follows a multi-agent architecture:

1. **Sources Layer**: CSV data files, OpenAI API
2. **Orchestration Layer**: Streamlit UI, Kosmos main orchestrator
3. **Processing Layer**: Data Analysis Agent, Literature Agent, Statistical Engine
4. **State & Storage Layer**: World Model, Local file storage
5. **Outputs Layer**: Analysis scripts, final reports

## Link Classification
- **Verified Links** (solid): All connections verified in source code
- **Inferred Links** (dotted): None - all flows traced through actual imports/calls

## Node Count: 11 (under 12 limit)
## Edge Count: 12 (under 16 limit)

## Assumptions
- Main entry point: streamlit_app_ULTIMATE.py (web) or main.py (CLI)
- Statistical analysis uses scipy (verified in imports)
- OpenAI API used for LLM features (optional, verified in code)
- All state persisted to local JSON files (verified)

## Unresolved Questions
None - all architecture components verified in codebase.
"""

    with open("READ_ME_ARCH_SLIDE.md", "w") as f:
        f.write(readme)
    print("✓ Saved: READ_ME_ARCH_SLIDE.md")

if __name__ == "__main__":
    # Generate graph data
    graph = create_architecture_graph()

    # Create PPTX
    output_file = create_presentation()

    # Validate
    if validate_pptx(output_file):
        # Save artifacts
        save_supporting_artifacts(graph)

        print(f"\n{'='*60}")
        print("✓ COMPLETE: Architecture overview slide generated")
        print(f"{'='*60}")
        print(f"Main file: {output_file}")
        print("Supporting files:")
        print("  - architecture_graph.json")
        print("  - component_index.csv")
        print("  - READ_ME_ARCH_SLIDE.md")
    else:
        print("\n✗ VALIDATION FAILED - check output")
        exit(1)
