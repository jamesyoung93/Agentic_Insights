#!/usr/bin/env python3
"""
Generate Architecture Overview PowerPoint
Consulting-style, diagram-first presentation
"""

import json
import csv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import matplotlib.lines as mlines

# Configuration
RUN_ID = "20251112_0348"
ARTIFACTS_DIR = Path("artifacts/agent_slide")
RUN_DIR = ARTIFACTS_DIR / RUN_ID
OUTPUT_FILE = "architecture_overview.pptx"
ACCENT_COLOR = "blue"
SLIDE_SIZE = "16x9"

def create_architecture_diagram():
    """Create architecture diagram as PNG using matplotlib"""

    # Load architecture graph
    with open(RUN_DIR / "architecture_graph.json", "r") as f:
        arch_data = json.load(f)

    # Create figure
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis('off')

    # Define swimlane positions
    lanes = {
        "Sources": {"y": 6.5, "height": 2, "color": "#FFE6E6"},
        "Processing": {"y": 3.5, "height": 2.5, "color": "#E8F4F8"},
        "Storage": {"y": 1.5, "height": 1.5, "color": "#FFF9E6"},
        "Outputs": {"y": 0.2, "height": 1, "color": "#E8F8F5"}
    }

    # Draw swimlanes
    for lane_name, lane_props in lanes.items():
        rect = FancyBboxPatch(
            (0.5, lane_props["y"] - lane_props["height"]),
            15,
            lane_props["height"],
            boxstyle="round,pad=0.05",
            linewidth=1.5,
            edgecolor='#95a5a6',
            facecolor=lane_props["color"],
            alpha=0.3,
            zorder=0
        )
        ax.add_patch(rect)
        ax.text(
            0.8,
            lane_props["y"] - 0.15,
            lane_name.upper(),
            fontsize=11,
            fontweight='bold',
            color='#34495e',
            verticalalignment='top'
        )

    # Node positions (manually placed for clarity)
    node_positions = {
        # Sources
        "n_csv_data": (3, 7.3),
        "n_user_input": (6, 7.3),
        "n_literature": (9, 7.3),

        # Processing
        "n_streamlit_app": (4.5, 5.2),
        "n_stats_engine": (7.5, 5.2),
        "n_llm_gateway": (10.5, 5.2),
        "n_world_model": (13.5, 5.2),

        # Storage
        "n_world_model_store": (11, 2.3),
        "n_session_store": (13.5, 2.3),

        # Outputs
        "n_dashboard": (11, 0.6),
        "n_reports": (13.5, 0.6),
    }

    # Node shapes and labels
    node_info = {
        "n_csv_data": {"label": "CSV Data", "shape": "cylinder", "color": "#C0392B"},
        "n_user_input": {"label": "User Input", "shape": "parallelogram", "color": "#E74C3C"},
        "n_literature": {"label": "Literature", "shape": "cylinder", "color": "#C0392B"},
        "n_streamlit_app": {"label": "Streamlit\nOrchestrator", "shape": "diamond", "color": "#3498DB"},
        "n_stats_engine": {"label": "Statistical\nAnalysis", "shape": "rect", "color": "#2980B9"},
        "n_llm_gateway": {"label": "LLM\nGateway", "shape": "rect", "color": "#2980B9"},
        "n_world_model": {"label": "World\nModel", "shape": "rect", "color": "#2980B9"},
        "n_world_model_store": {"label": "World Model\nJSON", "shape": "cylinder", "color": "#F39C12"},
        "n_session_store": {"label": "Session\nState", "shape": "cylinder", "color": "#F39C12"},
        "n_dashboard": {"label": "Dashboard", "shape": "parallelogram", "color": "#27AE60"},
        "n_reports": {"label": "Reports", "shape": "parallelogram", "color": "#27AE60"},
    }

    # Draw nodes
    for node_id, (x, y) in node_positions.items():
        info = node_info[node_id]

        if info["shape"] == "cylinder":
            # Draw cylinder (database)
            ellipse1 = patches.Ellipse((x, y + 0.15), 0.8, 0.15,
                                      facecolor='white', edgecolor=info["color"],
                                      linewidth=2, zorder=2)
            rect = patches.Rectangle((x - 0.4, y - 0.3), 0.8, 0.45,
                                    facecolor='white', edgecolor=info["color"],
                                    linewidth=2, zorder=2)
            ellipse2 = patches.Ellipse((x, y - 0.3), 0.8, 0.15,
                                      facecolor='white', edgecolor=info["color"],
                                      linewidth=2, zorder=2)
            ax.add_patch(rect)
            ax.add_patch(ellipse1)
            ax.add_patch(ellipse2)
        elif info["shape"] == "diamond":
            # Draw diamond (decision/orchestrator)
            diamond = patches.Polygon(
                [(x, y + 0.3), (x + 0.4, y), (x, y - 0.3), (x - 0.4, y)],
                facecolor='white', edgecolor=info["color"],
                linewidth=2.5, zorder=2
            )
            ax.add_patch(diamond)
        elif info["shape"] == "parallelogram":
            # Draw parallelogram (input/output)
            para = patches.Polygon(
                [(x - 0.3, y - 0.25), (x - 0.15, y + 0.25),
                 (x + 0.5, y + 0.25), (x + 0.35, y - 0.25)],
                facecolor='white', edgecolor=info["color"],
                linewidth=2, zorder=2
            )
            ax.add_patch(para)
        else:  # rectangle
            rect = FancyBboxPatch(
                (x - 0.4, y - 0.25), 0.8, 0.5,
                boxstyle="round,pad=0.05",
                facecolor='white', edgecolor=info["color"],
                linewidth=2, zorder=2
            )
            ax.add_patch(rect)

        # Add label
        ax.text(x, y, info["label"],
               ha='center', va='center',
               fontsize=8, fontweight='normal',
               color='#2c3e50', zorder=3)

    # Draw edges (arrows)
    edges = [
        ("n_csv_data", "n_streamlit_app", "CSV batch", "solid"),
        ("n_user_input", "n_streamlit_app", "params", "solid"),
        ("n_literature", "n_llm_gateway", "papers", "dotted"),
        ("n_streamlit_app", "n_stats_engine", "questions", "solid"),
        ("n_streamlit_app", "n_llm_gateway", "prompts", "solid"),
        ("n_stats_engine", "n_world_model", "results", "solid"),
        ("n_llm_gateway", "n_world_model", "discoveries", "solid"),
        ("n_world_model", "n_world_model_store", "JSON", "solid"),
        ("n_streamlit_app", "n_session_store", "state", "solid"),
        ("n_world_model_store", "n_reports", "data", "solid"),
        ("n_session_store", "n_dashboard", "live data", "solid"),
    ]

    for source, target, label, style in edges:
        if source in node_positions and target in node_positions:
            x1, y1 = node_positions[source]
            x2, y2 = node_positions[target]

            # Adjust start and end points to edge of boxes
            dx = x2 - x1
            dy = y2 - y1
            length = (dx**2 + dy**2)**0.5

            if length > 0:
                # Offset from center
                offset = 0.5
                x1_adj = x1 + (dx / length) * offset
                y1_adj = y1 + (dy / length) * offset
                x2_adj = x2 - (dx / length) * offset
                y2_adj = y2 - (dy / length) * offset

                linestyle = '--' if style == "dotted" else '-'

                arrow = FancyArrowPatch(
                    (x1_adj, y1_adj), (x2_adj, y2_adj),
                    arrowstyle='->',
                    linewidth=1.5 if style == "solid" else 1.0,
                    linestyle=linestyle,
                    color='#34495e',
                    mutation_scale=15,
                    zorder=1
                )
                ax.add_patch(arrow)

                # Add label
                mid_x = (x1 + x2) / 2
                mid_y = (y1 + y2) / 2
                ax.text(mid_x, mid_y, label,
                       ha='center', va='bottom',
                       fontsize=6, style='italic',
                       color='#7f8c8d',
                       bbox=dict(boxstyle='round,pad=0.3',
                                facecolor='white',
                                edgecolor='none',
                                alpha=0.8),
                       zorder=2)

    # Add legend
    legend_elements = [
        mlines.Line2D([], [], color='#34495e', linewidth=2,
                     linestyle='-', label='Verified'),
        mlines.Line2D([], [], color='#34495e', linewidth=1.5,
                     linestyle='--', label='Inferred'),
        patches.Patch(facecolor='white', edgecolor='#C0392B',
                     linewidth=2, label='Data Source'),
        patches.Patch(facecolor='white', edgecolor='#3498DB',
                     linewidth=2, label='Process'),
        patches.Patch(facecolor='white', edgecolor='#F39C12',
                     linewidth=2, label='Storage'),
        patches.Patch(facecolor='white', edgecolor='#27AE60',
                     linewidth=2, label='Output'),
    ]
    ax.legend(handles=legend_elements, loc='lower right',
             fontsize=7, frameon=True, fancybox=True)

    # Save figure
    diagram_path = RUN_DIR / "architecture_diagram.png"
    plt.tight_layout()
    plt.savefig(diagram_path, dpi=220, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    print(f"✅ Diagram saved: {diagram_path}")
    return diagram_path


def create_presentation():
    """Create PowerPoint presentation with python-pptx"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16) if SLIDE_SIZE == "16x9" else Inches(10)
    prs.slide_height = Inches(9) if SLIDE_SIZE == "16x9" else Inches(7.5)

    # === SLIDE 1: Architecture Overview ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Agentic Insights: System Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)
    title_para.alignment = PP_ALIGN.CENTER

    # Caption
    caption_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(15), Inches(0.4)
    )
    caption_frame = caption_box.text_frame
    caption_frame.text = "Autonomous AI discovery system with statistical rigor and LLM synthesis"
    caption_para = caption_frame.paragraphs[0]
    caption_para.font.size = Pt(16)
    caption_para.font.color.rgb = RGBColor(127, 140, 141)
    caption_para.alignment = PP_ALIGN.CENTER

    # Add diagram
    diagram_path = create_architecture_diagram()
    slide1.shapes.add_picture(
        str(diagram_path),
        Inches(0.5), Inches(1.5),
        width=Inches(15)
    )

    # Speaker notes
    notes_slide = slide1.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = """Architecture Overview:

SOURCES:
- CSV Data (data/customers.csv, data/competitor_data.csv): Customer transaction and competitor data
- User Input: Research objectives, API keys, and configuration via Streamlit UI
- Literature: Research papers stored in knowledge/literature/*.txt

PROCESSING:
- Streamlit Orchestrator: Main application controller that manages discovery cycles
- Statistical Analysis: scipy/numpy-based statistical tests (correlations, t-tests, ANOVA, regression)
- LLM Gateway: OpenAI API integration for question generation and discovery synthesis
- World Model: Knowledge graph manager tracking discoveries, trajectories, and hypotheses

STORAGE:
- World Model JSON: Persistent storage of discoveries and analysis trajectories
- Session State: Streamlit session state for runtime data preservation

OUTPUTS:
- Dashboard: Interactive Streamlit web interface with real-time visualizations
- Reports: Auto-generated text reports with statistical evidence extraction

Evidence: All components verified in source code (see component_index.csv)
"""

    # === SLIDE 2: Component Table (Appendix) ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.5)
    )
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Appendix A: Component Details"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(28)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(44, 62, 80)

    # Add component table
    # Read component index
    with open(RUN_DIR / "component_index.csv", "r") as f:
        reader = csv.DictReader(f)
        components = list(reader)

    # Group by lane
    lanes_components = {
        "Sources": [],
        "Processing": [],
        "Storage": [],
        "Outputs": []
    }

    for comp in components:
        comp_type = comp["type"]
        if comp_type == "source":
            lanes_components["Sources"].append(comp)
        elif comp_type == "process":
            lanes_components["Processing"].append(comp)
        elif comp_type == "store":
            lanes_components["Storage"].append(comp)
        elif comp_type == "output":
            lanes_components["Outputs"].append(comp)

    # Create text box with component details
    y_pos = 1.0
    for lane, comps in lanes_components.items():
        if comps:
            # Lane header
            lane_box = slide2.shapes.add_textbox(
                Inches(0.5), Inches(y_pos), Inches(15), Inches(0.3)
            )
            lane_frame = lane_box.text_frame
            lane_frame.text = f"{lane}:"
            lane_para = lane_frame.paragraphs[0]
            lane_para.font.size = Pt(14)
            lane_para.font.bold = True
            lane_para.font.color.rgb = RGBColor(52, 152, 219)
            y_pos += 0.35

            # Components
            for comp in comps[:3]:  # Limit to 3 per lane for space
                comp_box = slide2.shapes.add_textbox(
                    Inches(0.7), Inches(y_pos), Inches(14.5), Inches(0.4)
                )
                comp_frame = comp_box.text_frame
                comp_frame.word_wrap = True
                comp_frame.text = f"• {comp['name']}: {comp['inputs']} → {comp['outputs']}"
                comp_para = comp_frame.paragraphs[0]
                comp_para.font.size = Pt(10)
                comp_para.font.color.rgb = RGBColor(44, 62, 80)
                y_pos += 0.45

    # === SLIDE 3: Sequence Diagram (Appendix) ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box3 = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(15), Inches(0.5)
    )
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Appendix B: Discovery Cycle Sequence"
    title_para3 = title_frame3.paragraphs[0]
    title_para3.font.size = Pt(28)
    title_para3.font.bold = True
    title_para3.font.color.rgb = RGBColor(44, 62, 80)

    # Add sequence text
    seq_box = slide3.shapes.add_textbox(
        Inches(1), Inches(1.2), Inches(14), Inches(7)
    )
    seq_frame = seq_box.text_frame
    seq_frame.word_wrap = True

    sequence_text = """Discovery Cycle Flow:

1. User initiates discovery via Streamlit UI
   → Provides research objective and configures cycles

2. Data Loader reads CSV files
   → Loads and preprocesses customer/competitor data

3. Orchestrator generates research questions
   → Uses LLM (if enabled) or default question bank

4. Statistical Analysis Engine executes analyses
   → Performs correlations, t-tests, ANOVA, regression
   → Computes p-values, effect sizes, confidence intervals

5. LLM Gateway (optional) searches literature
   → Queries literature index for relevant papers
   → Extracts and synthesizes findings

6. Discovery Synthesis
   → LLM combines statistical + literature evidence
   → Generates actionable discoveries with confidence scores

7. World Model updates
   → Stores discoveries, trajectories, hypotheses
   → Persists to world_model.json

8. Report Generation
   → Auto-extracts statistics from trajectories
   → Formats enhanced report with full evidence

9. Dashboard displays results
   → Real-time visualization of discoveries
   → Interactive filtering and exploration
"""

    seq_frame.text = sequence_text
    for para in seq_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.color.rgb = RGBColor(44, 62, 80)
        para.space_before = Pt(6)

    return prs


def validate_presentation(prs_path):
    """Validate the generated presentation"""
    import zipfile

    print("\n🔍 Validating presentation...")

    # Check 1: Is it a valid zip file?
    if not zipfile.is_zipfile(prs_path):
        print("❌ VALIDATION FAILED: Not a valid ZIP file")
        return False
    print("✅ Valid ZIP structure")

    # Check 2: Can python-pptx reopen it?
    try:
        test_prs = Presentation(prs_path)
        slide_count = len(test_prs.slides)
        if slide_count < 1:
            print("❌ VALIDATION FAILED: No slides found")
            return False
        print(f"✅ Reopened successfully ({slide_count} slides)")
    except Exception as e:
        print(f"❌ VALIDATION FAILED: Cannot reopen - {e}")
        return False

    # Check 3: File size reasonable
    file_size = Path(prs_path).stat().st_size
    if file_size > 25 * 1024 * 1024:  # 25 MB limit
        print(f"⚠️ WARNING: File size ({file_size / 1024 / 1024:.1f} MB) exceeds 25 MB")
    else:
        print(f"✅ File size OK ({file_size / 1024 / 1024:.2f} MB)")

    print("✅ All validations passed!")
    return True


def main():
    """Main generation function"""
    print("="*80)
    print("ARCHITECTURE OVERVIEW POWERPOINT GENERATOR")
    print("="*80)
    print(f"RUN_ID: {RUN_ID}")
    print(f"Output: {OUTPUT_FILE}")
    print()

    # Create presentation
    print("📊 Creating presentation...")
    prs = create_presentation()

    # Save with atomic write
    output_path = Path(OUTPUT_FILE)
    temp_path = output_path.with_suffix('.tmp.pptx')

    print(f"💾 Saving to {output_path}...")
    prs.save(str(temp_path))

    # Atomic rename
    temp_path.replace(output_path)

    # Validate
    if validate_presentation(output_path):
        print(f"\n✅ SUCCESS: {output_path}")
        print(f"   Location: {output_path.absolute()}")
        return 0
    else:
        print(f"\n❌ VALIDATION FAILED")
        return 1


if __name__ == "__main__":
    exit(main())
