"""
Create Architecture Overview PowerPoint Presentation
Consulting-style slide deck for Agentic Insights system
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json

# Load architecture graph for reference
with open('architecture_graph.json', 'r') as f:
    arch_data = json.load(f)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Define color palette (McKinsey-style: navy, gray, white, accent blue)
NAVY = RGBColor(0, 32, 96)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(192, 192, 192)
ACCENT_BLUE = RGBColor(0, 112, 192)
WHITE = RGBColor(255, 255, 255)
BG_LIGHT = RGBColor(245, 245, 245)

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=DARK_GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = 'Calibri'

    return textbox

def add_shape_with_text(slide, shape_type, left, top, width, height, text,
                       fill_color, text_color=WHITE, font_size=11, bold=False):
    """Add a shape with text"""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    # Fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # Border
    shape.line.color.rgb = DARK_GRAY
    shape.line.width = Pt(1)

    # Text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
            run.font.name = 'Calibri'

    return shape

def add_connector(slide, x1, y1, x2, y2, label="", color=DARK_GRAY):
    """Add an arrow connector with optional label"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)

    # Add arrowhead
    connector.line.end_arrow_type = 2  # Arrow

    # Add label if provided
    if label:
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2 - 0.15
        add_text_box(slide, mid_x - 0.3, mid_y, 0.6, 0.2, label,
                    font_size=8, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    return connector

# ============================================================================
# SLIDE 1: EXECUTIVE ARCHITECTURE OVERVIEW
# ============================================================================

slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
add_text_box(slide, 0.5, 0.3, 12, 0.5,
             "System Architecture: How Agentic Insights Produces Autonomous Discoveries",
             font_size=24, bold=True, color=NAVY)

# Subtitle / Action Caption
add_text_box(slide, 0.5, 0.85, 12, 0.3,
             "Multi-cycle discovery engine combining statistical rigor with optional LLM intelligence to generate evidence-backed insights",
             font_size=12, color=DARK_GRAY)

# Main architecture flow (left to right: 4 domains)
y_base = 2.2

# DOMAIN 1: DATA SOURCES
domain1_x = 0.8
add_text_box(slide, domain1_x, y_base - 0.5, 2.2, 0.3, "DATA SOURCES",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# CSV Data (cylinder shape)
add_shape_with_text(slide, MSO_SHAPE.CAN, domain1_x + 0.1, y_base, 1.0, 0.7,
                   "CSV Files\n(customers,\ncompetitor)", LIGHT_GRAY, DARK_GRAY, 9)

# Literature KB (cylinder shape)
add_shape_with_text(slide, MSO_SHAPE.CAN, domain1_x + 1.2, y_base, 1.0, 0.7,
                   "Literature\nKnowledge\nBase", LIGHT_GRAY, DARK_GRAY, 9)

# OpenAI API (hexagon - external)
add_shape_with_text(slide, MSO_SHAPE.HEXAGON, domain1_x + 0.5, y_base + 1.1, 1.2, 0.5,
                   "OpenAI API\n(optional)", RGBColor(255, 240, 200), DARK_GRAY, 9)

# DOMAIN 2: PROCESSING / CORE ENGINE
domain2_x = 3.8
add_text_box(slide, domain2_x, y_base - 0.5, 3.0, 0.3, "PROCESSING ENGINE",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Streamlit Orchestrator
add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, domain2_x + 0.1, y_base - 0.15, 2.8, 0.5,
                   "Streamlit UI & Orchestration Engine", ACCENT_BLUE, WHITE, 10, True)

# Sub-modules (stacked)
modules_y = y_base + 0.55
module_height = 0.42

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, domain2_x + 0.2, modules_y, 1.2, module_height,
                   "Question\nGenerator", RGBColor(220, 230, 240), DARK_GRAY, 9)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, domain2_x + 1.5, modules_y, 1.2, module_height,
                   "Statistical\nAnalysis", RGBColor(220, 230, 240), DARK_GRAY, 9)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, domain2_x + 0.2, modules_y + 0.52, 1.2, module_height,
                   "Literature\nSearch", RGBColor(220, 230, 240), DARK_GRAY, 9)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, domain2_x + 1.5, modules_y + 0.52, 1.2, module_height,
                   "Discovery\nSynthesizer", RGBColor(220, 230, 240), DARK_GRAY, 9)

# DOMAIN 3: STORAGE / MODELS
domain3_x = 7.5
add_text_box(slide, domain3_x, y_base - 0.5, 2.2, 0.3, "STORAGE & STATE",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# World Model (cylinder)
add_shape_with_text(slide, MSO_SHAPE.CAN, domain3_x + 0.3, y_base + 0.2, 1.6, 0.8,
                   "World Model\n(discoveries,\ntrajectories,\nhypotheses)", NAVY, WHITE, 9)

# Report Generator (rectangle)
add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, domain3_x + 0.4, y_base + 1.2, 1.4, 0.4,
                   "Report\nGenerator", RGBColor(220, 230, 240), DARK_GRAY, 9)

# DOMAIN 4: OUTPUTS / SERVING
domain4_x = 10.4
add_text_box(slide, domain4_x, y_base - 0.5, 2.4, 0.3, "OUTPUTS",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Interactive Dashboard (parallelogram)
add_shape_with_text(slide, MSO_SHAPE.PARALLELOGRAM, domain4_x + 0.2, y_base, 2.0, 0.5,
                   "Interactive\nDashboard", RGBColor(220, 240, 220), DARK_GRAY, 9)

# Report File (parallelogram)
add_shape_with_text(slide, MSO_SHAPE.PARALLELOGRAM, domain4_x + 0.2, y_base + 0.65, 2.0, 0.5,
                   "Enhanced Report\n(auto_enhanced_report.txt)", RGBColor(220, 240, 220), DARK_GRAY, 9)

# World Model JSON (parallelogram)
add_shape_with_text(slide, MSO_SHAPE.PARALLELOGRAM, domain4_x + 0.2, y_base + 1.3, 2.0, 0.5,
                   "World Model State\n(world_model.json)", RGBColor(220, 240, 220), DARK_GRAY, 9)

# ============================================================================
# ARROWS / DATA FLOW
# ============================================================================

# CSV → Processing
add_connector(slide, domain1_x + 1.1, y_base + 0.35, domain2_x + 0.1, y_base + 0.1, "CSV")

# Literature KB → Lit Search
add_connector(slide, domain1_x + 2.2, y_base + 0.35, domain2_x + 0.3, modules_y + 0.73, "papers")

# OpenAI → Question Gen
add_connector(slide, domain1_x + 1.7, y_base + 1.35, domain2_x + 0.4, modules_y + 0.2, "REST/JSON")

# OpenAI → Lit Search (dotted - inferred)
add_connector(slide, domain1_x + 1.5, y_base + 1.5, domain2_x + 0.5, modules_y + 0.9, "REST/JSON")

# OpenAI → Synthesizer
add_connector(slide, domain1_x + 1.7, y_base + 1.45, domain2_x + 1.7, modules_y + 1.15, "REST/JSON")

# Processing → World Model
add_connector(slide, domain2_x + 2.9, y_base + 0.6, domain3_x + 0.3, y_base + 0.5,
             "discoveries\n+ trajectories")

# World Model → Report Gen
add_connector(slide, domain3_x + 1.1, y_base + 1.0, domain3_x + 1.1, y_base + 1.2)

# Report Gen → Report File
add_connector(slide, domain3_x + 1.8, y_base + 1.4, domain4_x + 0.2, y_base + 0.9, "write TXT")

# World Model → JSON Output
add_connector(slide, domain3_x + 1.9, y_base + 0.8, domain4_x + 0.2, y_base + 1.55, "save JSON")

# Streamlit → Dashboard
add_connector(slide, domain2_x + 2.9, y_base + 0.1, domain4_x + 0.2, y_base + 0.25, "HTTP/render")

# ============================================================================
# LEGEND
# ============================================================================

legend_x = 0.5
legend_y = 6.2

add_text_box(slide, legend_x, legend_y, 2, 0.2, "LEGEND",
             font_size=9, bold=True, color=DARK_GRAY)

# Shapes legend
add_shape_with_text(slide, MSO_SHAPE.CAN, legend_x, legend_y + 0.25, 0.4, 0.3,
                   "", LIGHT_GRAY, DARK_GRAY, 8)
add_text_box(slide, legend_x + 0.45, legend_y + 0.27, 1.2, 0.2, "= Data Store",
             font_size=8, color=DARK_GRAY)

add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, legend_x + 1.8, legend_y + 0.25, 0.4, 0.3,
                   "", RGBColor(220, 230, 240), DARK_GRAY, 8)
add_text_box(slide, legend_x + 2.25, legend_y + 0.27, 1.5, 0.2, "= Processing Module",
             font_size=8, color=DARK_GRAY)

add_shape_with_text(slide, MSO_SHAPE.HEXAGON, legend_x + 4.0, legend_y + 0.25, 0.4, 0.3,
                   "", RGBColor(255, 240, 200), DARK_GRAY, 8)
add_text_box(slide, legend_x + 4.45, legend_y + 0.27, 1.5, 0.2, "= External System",
             font_size=8, color=DARK_GRAY)

add_shape_with_text(slide, MSO_SHAPE.PARALLELOGRAM, legend_x + 6.0, legend_y + 0.25, 0.4, 0.3,
                   "", RGBColor(220, 240, 220), DARK_GRAY, 8)
add_text_box(slide, legend_x + 6.45, legend_y + 0.27, 1.2, 0.2, "= Output Artifact",
             font_size=8, color=DARK_GRAY)

# Runtime trigger note
add_text_box(slide, 8.5, legend_y, 4.3, 0.6,
             "RUNTIME TRIGGER: User initiates discovery → System runs 1-20 cycles → Each cycle:\n" +
             "1) Generates questions  2) Runs statistical tests  3) Searches literature (optional)\n" +
             "4) Synthesizes discoveries  5) Updates world model  6) Produces reports",
             font_size=8, color=DARK_GRAY)

# ============================================================================
# SLIDE 2: COMPONENT DETAIL TABLE
# ============================================================================

slide2 = prs.slides.add_slide(slide_layout)
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = WHITE

add_text_box(slide2, 0.5, 0.3, 12, 0.4, "Appendix: Component Details",
             font_size=22, bold=True, color=NAVY)

# Table header
table_y = 1.0
col_widths = [2.5, 2.2, 2.8, 2.8, 2.2]
headers = ["Component", "Path", "Key Functions", "Inputs", "Outputs"]

for i, (header, width) in enumerate(zip(headers, col_widths)):
    x = 0.5 + sum(col_widths[:i])
    add_text_box(slide2, x, table_y, width, 0.3, header,
                font_size=10, bold=True, color=WHITE)
    # Background for header
    header_bg = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(table_y), Inches(width), Inches(0.3)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = NAVY
    header_bg.line.color.rgb = WHITE
    # Move text to front
    slide2.shapes._spTree.remove(header_bg._element)
    slide2.shapes._spTree.insert(2, header_bg._element)

# Table rows (sample - key components only)
rows = [
    ["Streamlit UI", "streamlit_app_\nULTIMATE.py", "main(), run_discovery_\ncycle()",
     "User inputs,\nCSV files", "Dashboard,\nlogs"],
    ["World Model", "world_model_\nbuilder.py", "add_discovery(),\nadd_trajectory()",
     "Discovery/\ntrajectory data", "world_model.\njson"],
    ["Report Gen", "auto_enhanced_\nreport.py", "generate_enhanced_\nreport()",
     "Discoveries,\ntrajectories", "auto_enhanced_\nreport.txt"],
    ["Stats Engine", "embedded in\nstreamlit", "perform_statistical_\nanalysis()",
     "DataFrame,\nquestion", "p-values,\neffect sizes"],
    ["Lit Search", "agents/\nliterature_\nsearcher.py", "search(),\n_extract_insights()",
     "Query,\nliterature KB", "Papers,\ninsights"],
]

row_y = table_y + 0.35
for row_data in rows:
    for i, (cell, width) in enumerate(zip(row_data, col_widths)):
        x = 0.5 + sum(col_widths[:i])
        # Cell background
        cell_bg = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(row_y), Inches(width), Inches(0.6)
        )
        cell_bg.fill.solid()
        cell_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
        cell_bg.line.color.rgb = LIGHT_GRAY

        # Cell text
        add_text_box(slide2, x + 0.05, row_y + 0.05, width - 0.1, 0.5, cell,
                    font_size=8, color=DARK_GRAY)

    row_y += 0.65

# Dependencies note
add_text_box(slide2, 0.5, row_y + 0.2, 12, 0.8,
             "CRITICAL DEPENDENCIES:\n" +
             "• Python 3.8+  • Streamlit (UI framework)  • SciPy/NumPy (statistical analysis)  • Pandas (data manipulation)\n" +
             "• OpenAI API (optional, for LLM features)  • python-pptx (presentation generation)\n\n" +
             "DEPLOYMENT: Runs locally or containerized (Docker). No cloud infrastructure required for core functionality.",
             font_size=9, color=DARK_GRAY)

# ============================================================================
# SLIDE 3: PRIMARY USER JOURNEY
# ============================================================================

slide3 = prs.slides.add_slide(slide_layout)
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = WHITE

add_text_box(slide3, 0.5, 0.3, 12, 0.4, "Appendix: Primary User Journey",
             font_size=22, bold=True, color=NAVY)

add_text_box(slide3, 0.5, 0.8, 12, 0.3,
             "End-to-end flow: User initiates discovery → System autonomously generates insights → Reports delivered",
             font_size=11, color=DARK_GRAY)

# Sequence diagram (simplified vertical flow)
seq_x = 2.0
seq_y = 1.5
step_height = 0.7

steps = [
    ("1. User Setup", "Configure objective, API key (optional), # cycles", ACCENT_BLUE),
    ("2. Data Load", "Load CSV files or generate sample data", RGBColor(100, 150, 200)),
    ("3. Discovery Loop", "For each cycle (1-20):", RGBColor(80, 130, 180)),
    ("   3a. Generate Questions", "LLM or predefined research questions", RGBColor(180, 200, 220)),
    ("   3b. Run Statistical Tests", "Correlations, t-tests, ANOVA, regression", RGBColor(180, 200, 220)),
    ("   3c. Search Literature", "(Optional) Query knowledge base via LLM", RGBColor(180, 200, 220)),
    ("   3d. Synthesize", "Combine stats + literature → discoveries", RGBColor(180, 200, 220)),
    ("   3e. Update World Model", "Add discoveries & trajectories to graph", RGBColor(180, 200, 220)),
    ("4. Generate Report", "Extract statistics, format enhanced report", RGBColor(100, 150, 200)),
    ("5. Deliver Outputs", "Dashboard + TXT report + JSON model", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(steps):
    y = seq_y + i * step_height

    # Step box
    add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, seq_x, y, 4.0, 0.5,
                       title, color, WHITE, 10, True)

    # Description
    add_text_box(slide3, seq_x + 4.2, y + 0.05, 6.0, 0.4, desc,
                font_size=9, color=DARK_GRAY)

    # Arrow to next step (except last)
    if i < len(steps) - 1:
        add_connector(slide3, seq_x + 2.0, y + 0.5, seq_x + 2.0, y + 0.7, "")

# Timing estimate
add_text_box(slide3, seq_x + 7.0, seq_y + len(steps) * step_height - 0.5, 4.0, 0.8,
             "TYPICAL TIMING:\n" +
             "• 5 cycles with LLM: ~10-20 min\n" +
             "• 5 cycles statistical-only: ~2-5 min\n" +
             "• Scales with data size & API latency",
             font_size=9, color=DARK_GRAY)

# ============================================================================
# SLIDE 4: DEPLOYMENT VIEW
# ============================================================================

slide4 = prs.slides.add_slide(slide_layout)
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = WHITE

add_text_box(slide4, 0.5, 0.3, 12, 0.4, "Appendix: Deployment Architecture",
             font_size=22, bold=True, color=NAVY)

# Three deployment tiers
tier_y = 1.5
tier_width = 3.5
tier_gap = 0.5

# Local Development
add_text_box(slide4, 0.8, tier_y - 0.3, tier_width, 0.2, "LOCAL DEVELOPMENT",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_shape_with_text(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, tier_y, tier_width, 1.2,
                   "Python Environment\n\n" +
                   "• Streamlit dev server\n" +
                   "• Local file storage\n" +
                   "• Direct OpenAI calls\n" +
                   "• No containerization",
                   RGBColor(230, 240, 250), DARK_GRAY, 9)

add_text_box(slide4, 0.8, tier_y + 1.3, tier_width, 0.3, "Run: streamlit run\nstreamlit_app_ULTIMATE.py",
             font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Containerized
add_text_box(slide4, 0.8 + tier_width + tier_gap, tier_y - 0.3, tier_width, 0.2,
             "CONTAINERIZED",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_shape_with_text(slide4, MSO_SHAPE.ROUNDED_RECTANGLE,
                   0.8 + tier_width + tier_gap, tier_y, tier_width, 1.2,
                   "Docker Container\n\n" +
                   "• Dockerfile included\n" +
                   "• Isolated dependencies\n" +
                   "• Volume mounts for data\n" +
                   "• Port 8501 exposed",
                   RGBColor(230, 240, 250), DARK_GRAY, 9)

add_text_box(slide4, 0.8 + tier_width + tier_gap, tier_y + 1.3, tier_width, 0.3,
             "Run: docker build & run\n(see Dockerfile if exists)",
             font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Cloud (inferred/optional)
add_text_box(slide4, 0.8 + 2*(tier_width + tier_gap), tier_y - 0.3, tier_width, 0.2,
             "CLOUD (Optional)",
             font_size=10, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_shape_with_text(slide4, MSO_SHAPE.ROUNDED_RECTANGLE,
                   0.8 + 2*(tier_width + tier_gap), tier_y, tier_width, 1.2,
                   "Cloud Deployment\n\n" +
                   "• AWS/GCP/Azure VM\n" +
                   "• S3/GCS for data\n" +
                   "• Secrets manager for keys\n" +
                   "• Load balancer (if scaled)",
                   RGBColor(230, 240, 250), DARK_GRAY, 9)

add_text_box(slide4, 0.8 + 2*(tier_width + tier_gap), tier_y + 1.3, tier_width, 0.3,
             "Deploy: Cloud-specific\n(Terraform/Helm if infra exists)",
             font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Environment variables / config
add_text_box(slide4, 0.8, tier_y + 2.0, 12, 1.0,
             "CONFIGURATION & SECRETS:\n\n" +
             "• OpenAI API Key: Set in config.py or via Streamlit UI (saved to session)\n" +
             "• Data Sources: CSV files in data/ directory or auto-generated\n" +
             "• Literature KB: JSON index + text files in knowledge/literature/\n" +
             "• Outputs: world_model.json & auto_enhanced_report.txt in base directory\n\n" +
             "SECURITY NOTE: API keys are stored in config.py (local file). For production, use environment variables or secrets manager.",
             font_size=9, color=DARK_GRAY)

# ============================================================================
# SAVE PRESENTATION
# ============================================================================

prs.save('architecture_overview.pptx')
print("✅ PowerPoint presentation created: architecture_overview.pptx")
print("   - Slide 1: Executive Architecture Overview (main flow diagram)")
print("   - Slide 2: Component Details (table)")
print("   - Slide 3: Primary User Journey (sequence)")
print("   - Slide 4: Deployment Architecture (environment tiers)")
