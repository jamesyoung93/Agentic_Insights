"""
Agentic Insights PowerPoint Presentation Generator

This script generates a comprehensive PowerPoint presentation explaining the
Agentic Insights system architecture, concepts, and capabilities.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    """Add a title slide to the presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, content_items, layout_idx=1):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title

    content = slide.placeholders[1].text_frame
    content.clear()

    for item in content_items:
        if isinstance(item, dict):
            p = content.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            if item.get('bold'):
                p.font.bold = True
        else:
            p = content.add_paragraph()
            p.text = item
            p.level = 0

    return slide


def add_section_header(prs, title, subtitle=""):
    """Add a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle
        except:
            pass
    return slide


def create_presentation():
    """Create the complete Agentic Insights presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Agentic Insights",
        "Autonomous Data-Driven Scientific Discovery System"
    )

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "AI-powered autonomous research system for scientific discovery",
        "Combines multi-agent orchestration with rigorous statistical analysis",
        "Inspired by the Kosmos framework (arXiv:2511.02824v2)",
        "Automatically explores datasets and generates research-grade insights",
        "Iterative learning approach that builds knowledge across cycles",
        "Dual evidence: Data analysis + Literature review synthesis"
    ])

    # Slide 3: What Problem Does It Solve?
    add_content_slide(prs, "What Problem Does It Solve?", [
        {"text": "Manual Bottlenecks in Research", "bold": True},
        {"text": "Traditional data analysis requires extensive manual effort", "level": 1},
        {"text": "Literature review is time-consuming and repetitive", "level": 1},
        {"text": "Hypothesis generation limited by human cognitive constraints", "level": 1},
        "",
        {"text": "Quality & Reproducibility Issues", "bold": True},
        {"text": "Inconsistent statistical rigor across analyses", "level": 1},
        {"text": "Limited traceability from data to conclusions", "level": 1},
        {"text": "Difficulty maintaining context across multiple investigations", "level": 1}
    ])

    # Slide 4: System Architecture Overview
    add_section_header(prs, "System Architecture", "Multi-Layer Intelligent Design")

    # Slide 5: Architecture Layers
    add_content_slide(prs, "Architecture Layers", [
        {"text": "Orchestration Layer", "bold": True},
        {"text": "KosmosFramework coordinates multi-cycle discovery processes", "level": 1},
        "",
        {"text": "Agent Layer", "bold": True},
        {"text": "DataAnalysisAgent: Generates & executes statistical code", "level": 1},
        {"text": "LiteratureSearchAgent: Finds and synthesizes research", "level": 1},
        "",
        {"text": "Knowledge Layer", "bold": True},
        {"text": "World Model: Persistent context across discovery cycles", "level": 1},
        "",
        {"text": "Enhancement Layer", "bold": True},
        {"text": "Statistical extraction and scientific report generation", "level": 1},
        "",
        {"text": "Interface Layer", "bold": True},
        {"text": "Streamlit web applications for user interaction", "level": 1}
    ])

    # Slide 6: Core Concepts
    add_section_header(prs, "Core Concepts", "Autonomous Discovery Methodology")

    # Slide 7: Autonomous Discovery Cycle
    add_content_slide(prs, "Autonomous Discovery Cycle", [
        {"text": "Step 1: Review Current Knowledge State", "bold": True},
        {"text": "World Model provides context from previous discoveries", "level": 1},
        "",
        {"text": "Step 2: Generate Research Questions", "bold": True},
        {"text": "LLM creates novel hypotheses based on existing knowledge", "level": 1},
        "",
        {"text": "Step 3: Parallel Task Execution", "bold": True},
        {"text": "Data analysis + Literature search run simultaneously", "level": 1},
        "",
        {"text": "Step 4: Synthesize Discoveries", "bold": True},
        {"text": "Combine evidence from multiple sources", "level": 1},
        "",
        {"text": "Step 5: Update World Model", "bold": True},
        {"text": "Persist findings for next cycle", "level": 1}
    ])

    # Slide 8: World Model Concept
    add_content_slide(prs, "The World Model Concept", [
        {"text": "Persistent Knowledge Repository", "bold": True},
        {"text": "Structured representation of accumulated knowledge", "level": 1},
        {"text": "Tracks discoveries, trajectories, hypotheses, questions", "level": 1},
        "",
        {"text": "Context Management", "bold": True},
        {"text": "Generates summaries for agent prompts", "level": 1},
        {"text": "Enables cumulative learning across cycles", "level": 1},
        {"text": "Builds on previous findings systematically", "level": 1},
        "",
        {"text": "Enhanced Features", "bold": True},
        {"text": "Rich metadata: timestamps, confidence scores", "level": 1},
        {"text": "Advanced filtering and querying capabilities", "level": 1},
        {"text": "Cycle summaries and question management", "level": 1}
    ])

    # Slide 9: Scientific Rigor
    add_content_slide(prs, "Scientific Rigor Methodology", [
        {"text": "Automatic Statistical Extraction", "bold": True},
        {"text": "p-values, effect sizes, confidence intervals", "level": 1},
        {"text": "Pearson correlation, t-tests, ANOVA, linear regression", "level": 1},
        "",
        {"text": "Causal Inference Assessment", "bold": True},
        {"text": "Bradford Hill criteria for causality", "level": 1},
        {"text": "Explicit discussion of confounders", "level": 1},
        "",
        {"text": "Transparency & Traceability", "bold": True},
        {"text": "Every discovery linked to supporting trajectories", "level": 1},
        {"text": "Analysis code saved with unique IDs", "level": 1},
        {"text": "Complete audit trail from raw data to conclusions", "level": 1},
        {"text": "Methodology documentation included", "level": 1}
    ])

    # Slide 10: System Components
    add_section_header(prs, "System Components", "Key Modules and Their Functions")

    # Slide 11: Core Components Detail
    add_content_slide(prs, "Core Components", [
        {"text": "KosmosFramework (main.py - 291 lines)", "bold": True},
        {"text": "Central orchestrator, manages 5-10 discovery cycles", "level": 1},
        {"text": "Research question generation using GPT-3.5-turbo", "level": 1},
        "",
        {"text": "DataAnalysisAgent (agents/data_analyst.py - 237 lines)", "bold": True},
        {"text": "Generates executable Python code for analysis", "level": 1},
        {"text": "Safe execution with error handling (max 3 retries)", "level": 1},
        "",
        {"text": "LiteratureSearchAgent (agents/literature_searcher.py - 228 lines)", "bold": True},
        {"text": "Searches and ranks research papers by relevance", "level": 1},
        {"text": "Links literature support to data discoveries", "level": 1}
    ])

    # Slide 12: World Model & Report Generation
    add_content_slide(prs, "World Model & Report Generation", [
        {"text": "World Model (world_model_builder.py - 396 lines)", "bold": True},
        {"text": "Structured dataclasses for Discovery & Trajectory objects", "level": 1},
        {"text": "Rich metadata, filtering, and querying capabilities", "level": 1},
        {"text": "Context generation optimized for LLM prompts", "level": 1},
        "",
        {"text": "AutoEnhancedReportGenerator (auto_enhanced_report.py - 332 lines)", "bold": True},
        {"text": "Extracts statistics from trajectory outputs", "level": 1},
        {"text": "Formats discoveries with supporting evidence", "level": 1},
        {"text": "Includes methodology, limitations, future directions", "level": 1},
        "",
        {"text": "Streamlit Interface (streamlit_app_ULTIMATE.py - 1,642 lines)", "bold": True},
        {"text": "Real OpenAI API integration with live statistical analysis", "level": 1},
        {"text": "Progress tracking and comprehensive logging", "level": 1}
    ])

    # Slide 13: Component Interaction Flow
    add_content_slide(prs, "Component Interaction Flow", [
        "User Input (Streamlit)",
        {"text": "↓", "level": 0},
        "KosmosFramework Initialization",
        {"text": "↓", "level": 0},
        "Generate Research Questions (LLM)",
        {"text": "↓", "level": 0},
        "Parallel Execution:",
        {"text": "DataAnalysisAgent → Execute Code → Extract Stats", "level": 1},
        {"text": "LiteratureSearchAgent → Find Papers → Synthesize", "level": 1},
        {"text": "↓", "level": 0},
        "Update World Model",
        {"text": "↓", "level": 0},
        "Synthesize Discoveries (LLM)",
        {"text": "↓", "level": 0},
        "AutoEnhancedReportGenerator → Enhanced Scientific Report"
    ])

    # Slide 14: What It Accomplishes
    add_section_header(prs, "What It Accomplishes", "System Objectives and Outcomes")

    # Slide 15: Primary Objectives
    add_content_slide(prs, "Primary Objectives", [
        {"text": "Automated Scientific Discovery Pipeline", "bold": True},
        {"text": "Raw data → rigorous, publication-ready insights", "level": 1},
        {"text": "Minimal human intervention required", "level": 1},
        "",
        {"text": "Business Intelligence", "bold": True},
        {"text": "Identify key drivers of customer loyalty and revenue", "level": 1},
        {"text": "Investigate customer behavior relationships", "level": 1},
        {"text": "Analyze competitor impacts on customer behavior", "level": 1},
        {"text": "Optimize operational metrics (wait times, satisfaction)", "level": 1},
        "",
        {"text": "Scientific Quality Assurance", "bold": True},
        {"text": "Ensure proper statistical testing (p < 0.05)", "level": 1},
        {"text": "Eliminate superficial correlational claims", "level": 1},
        {"text": "Maintain complete traceability", "level": 1}
    ])

    # Slide 16: Use Cases
    add_section_header(prs, "Use Cases", "Applications Across Domains")

    # Slide 17: Current Domain - Coffee Shop Analytics
    add_content_slide(prs, "Current Domain: Coffee Shop Chain Analytics", [
        {"text": "Customer Analytics", "bold": True},
        {"text": "Customer segmentation and lifetime value analysis", "level": 1},
        {"text": "Loyalty program effectiveness measurement", "level": 1},
        {"text": "Mobile app impact on retention", "level": 1},
        "",
        {"text": "Operational Optimization", "bold": True},
        {"text": "Wait time optimization", "level": 1},
        {"text": "Geographic and seasonal pattern analysis", "level": 1},
        {"text": "Price sensitivity analysis", "level": 1},
        "",
        {"text": "Competitive Intelligence", "bold": True},
        {"text": "Competitor activity influence on behavior", "level": 1},
        {"text": "Market dynamics and customer churn prediction", "level": 1}
    ])

    # Slide 18: Generalizable Applications
    add_content_slide(prs, "Generalizable Applications", [
        {"text": "Healthcare", "bold": True},
        {"text": "Patient outcome analysis, treatment effectiveness studies", "level": 1},
        "",
        {"text": "Finance", "bold": True},
        {"text": "Investment pattern discovery, risk factor analysis", "level": 1},
        "",
        {"text": "E-commerce", "bold": True},
        {"text": "User behavior analysis, conversion optimization", "level": 1},
        "",
        {"text": "Academic Research", "bold": True},
        {"text": "Automated literature review + data synthesis", "level": 1},
        "",
        {"text": "Manufacturing", "bold": True},
        {"text": "Quality control, process optimization", "level": 1},
        "",
        {"text": "Marketing", "bold": True},
        {"text": "Campaign effectiveness, customer segmentation", "level": 1}
    ])

    # Slide 19: Technology Stack
    add_section_header(prs, "Technology Stack", "Modern AI and Data Science Tools")

    # Slide 20: Technology Stack Details
    add_content_slide(prs, "Technology Stack", [
        {"text": "AI/ML Layer", "bold": True},
        {"text": "OpenAI GPT-3.5-turbo for question generation and synthesis", "level": 1},
        {"text": "LLM-driven code generation with self-correction", "level": 1},
        "",
        {"text": "Data Science Stack", "bold": True},
        {"text": "pandas (>=2.0.0), numpy (>=1.24.0)", "level": 1},
        {"text": "scipy (>=1.10.0), scikit-learn (>=1.3.0)", "level": 1},
        {"text": "statsmodels (>=0.14.0), pingouin (>=0.5.3)", "level": 1},
        "",
        {"text": "Visualization", "bold": True},
        {"text": "matplotlib (>=3.7.0), seaborn (>=0.12.0)", "level": 1},
        "",
        {"text": "User Interface", "bold": True},
        {"text": "Streamlit (>=1.28.0) - Web-based interactive interface", "level": 1}
    ])

    # Slide 21: Implementation Features
    add_content_slide(prs, "Implementation Features", [
        {"text": "Safety & Reliability", "bold": True},
        {"text": "Safe code execution with retry logic (3 attempts)", "level": 1},
        {"text": "Data sanitization for statistical tests", "level": 1},
        {"text": "Timeout protection (300 seconds per analysis)", "level": 1},
        "",
        {"text": "Performance", "bold": True},
        {"text": "Handles datasets up to ~5GB", "level": 1},
        {"text": "Generates reports in < 1 second", "level": 1},
        {"text": "15-minute cache for web fetching", "level": 1},
        "",
        {"text": "State Management", "bold": True},
        {"text": "Persistent world model state across sessions", "level": 1},
        {"text": "Resume capability from any cycle", "level": 1},
        {"text": "Complete history tracking", "level": 1}
    ])

    # Slide 22: Key Differentiators
    add_section_header(prs, "Key Differentiators", "What Makes Agentic Insights Unique")

    # Slide 23: Differentiators Detail
    add_content_slide(prs, "What Sets Us Apart", [
        {"text": "1. True Autonomy", "bold": True},
        {"text": "Runs 10-20+ cycles without human intervention", "level": 1},
        "",
        {"text": "2. Scientific Rigor", "bold": True},
        {"text": "Automatic statistical extraction, not just narratives", "level": 1},
        "",
        {"text": "3. Dual Evidence Approach", "bold": True},
        {"text": "Combines data analysis + literature review", "level": 1},
        "",
        {"text": "4. Iterative Learning", "bold": True},
        {"text": "Each cycle builds on previous discoveries", "level": 1},
        "",
        {"text": "5. Complete Traceability", "bold": True},
        {"text": "Full audit trail from data to insights", "level": 1},
        "",
        {"text": "6. Production Ready", "bold": True},
        {"text": "Web UI, error handling, logging, state management", "level": 1}
    ])

    # Slide 24: System Workflow Example
    add_content_slide(prs, "System Workflow Example", [
        {"text": "Cycle 1: Initial Exploration", "bold": True},
        {"text": "Question: What factors influence customer retention?", "level": 1},
        {"text": "Analysis: Correlation between visits and loyalty status", "level": 1},
        {"text": "Discovery: Loyalty members visit 2.3x more often (p<0.001)", "level": 1},
        "",
        {"text": "Cycle 2: Building on Discovery", "bold": True},
        {"text": "Question: Does mobile app usage affect loyalty?", "level": 1},
        {"text": "Analysis: App users vs non-app users comparison", "level": 1},
        {"text": "Discovery: App users have 35% higher retention", "level": 1},
        "",
        {"text": "Cycle 3: Deeper Investigation", "bold": True},
        {"text": "Question: What drives app adoption?", "level": 1},
        {"text": "Analysis: Demographics and behavior patterns", "level": 1},
        {"text": "Discovery: Age < 40 and urban location predict adoption", "level": 1}
    ])

    # Slide 25: Codebase Statistics
    add_content_slide(prs, "Codebase Statistics", [
        {"text": "Scale & Complexity", "bold": True},
        {"text": "~25 Python files totaling 6,700+ lines of code", "level": 1},
        {"text": "4 specialized agents + orchestrator + UI", "level": 1},
        {"text": "4 comprehensive markdown documentation guides", "level": 1},
        "",
        {"text": "Data & Knowledge", "bold": True},
        {"text": "Sample dataset: 4,992 customer records with 9+ attributes", "level": 1},
        {"text": "Literature database: 8 research papers with metadata", "level": 1},
        "",
        {"text": "Code Organization", "bold": True},
        {"text": "agents/ - Specialized agent modules", "level": 1},
        {"text": "data/ - CSV datasets and generators", "level": 1},
        {"text": "knowledge/ - Literature database", "level": 1},
        {"text": "outputs/ - Generated analyses and reports", "level": 1}
    ])

    # Slide 26: Configuration & Customization
    add_content_slide(prs, "Configuration & Customization", [
        {"text": "Highly Configurable System", "bold": True},
        {"text": "config.py for central configuration", "level": 1},
        {"text": "Customizable research objectives", "level": 1},
        {"text": "Adjustable cycle counts and parallel tasks", "level": 1},
        {"text": "Temperature settings for LLM creativity", "level": 1},
        "",
        {"text": "Multi-Project Support", "bold": True},
        {"text": "Multiple base directory support", "level": 1},
        {"text": "Project-specific configurations", "level": 1},
        "",
        {"text": "Extensible Architecture", "bold": True},
        {"text": "Easy to add new agents", "level": 1},
        {"text": "Custom statistical tests can be integrated", "level": 1},
        {"text": "Pluggable literature sources", "level": 1}
    ])

    # Slide 27: Future Directions
    add_content_slide(prs, "Future Directions", [
        {"text": "Enhanced Capabilities", "bold": True},
        {"text": "Integration with more data sources", "level": 1},
        {"text": "Advanced causal inference methods", "level": 1},
        {"text": "Real-time data processing", "level": 1},
        "",
        {"text": "Expanded Applications", "bold": True},
        {"text": "Domain-specific agent libraries", "level": 1},
        {"text": "Multi-modal data analysis (text, images, time-series)", "level": 1},
        {"text": "Collaborative multi-agent scenarios", "level": 1},
        "",
        {"text": "Performance & Scale", "bold": True},
        {"text": "Distributed computing for larger datasets", "level": 1},
        {"text": "GPU acceleration for statistical computations", "level": 1},
        {"text": "Enhanced caching and incremental processing", "level": 1}
    ])

    # Slide 28: Getting Started
    add_content_slide(prs, "Getting Started", [
        {"text": "Installation", "bold": True},
        {"text": "pip install -r requirements.txt", "level": 1},
        {"text": "Configure OpenAI API key", "level": 1},
        "",
        {"text": "Quick Start", "bold": True},
        {"text": "streamlit run streamlit_app_ULTIMATE.py", "level": 1},
        {"text": "Load your dataset", "level": 1},
        {"text": "Define research objectives", "level": 1},
        {"text": "Run autonomous discovery", "level": 1},
        "",
        {"text": "Documentation", "bold": True},
        {"text": "README.md - Overview and setup", "level": 1},
        {"text": "docs/ - Detailed architecture guides", "level": 1},
        {"text": "examples/ - Sample analyses and outputs", "level": 1}
    ])

    # Slide 29: Conclusion
    add_content_slide(prs, "Conclusion", [
        {"text": "Revolutionary Approach to Data Science", "bold": True},
        {"text": "Combines autonomy with scientific rigor", "level": 1},
        {"text": "Reduces time from hypothesis to insight from weeks to hours", "level": 1},
        "",
        {"text": "Production-Ready System", "bold": True},
        {"text": "Comprehensive error handling and logging", "level": 1},
        {"text": "User-friendly interface with Streamlit", "level": 1},
        {"text": "Scalable architecture for various domains", "level": 1},
        "",
        {"text": "Continuous Innovation", "bold": True},
        {"text": "Based on latest research (Kosmos framework)", "level": 1},
        {"text": "Active development and improvements", "level": 1},
        {"text": "Extensible for future capabilities", "level": 1}
    ])

    # Slide 30: Thank You / Contact
    add_title_slide(
        prs,
        "Agentic Insights",
        "Questions? Explore the codebase and documentation for more details."
    )

    return prs


def main():
    """Main function to generate and save the presentation"""
    print("Generating Agentic Insights presentation...")

    prs = create_presentation()

    output_file = "Agentic_Insights_Presentation.pptx"
    prs.save(output_file)

    print(f"Presentation successfully created: {output_file}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
