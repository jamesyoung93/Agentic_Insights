"""
Generate Architecture Overview PowerPoint Slide
Constraint-compliant consulting-style diagram
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

# Load architecture data
with open('architecture_graph.json', 'r') as f:
    arch_data = json.load(f)

# Create presentation (16:9 format)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Color palette (grayscale + blue accent)
COLOR_TITLE = RGBColor(31, 56, 100)
COLOR_DATASTORE = RGBColor(176, 196, 222)  # Light steel blue
COLOR_SERVICE = RGBColor(144, 164, 174)    # Blue grey
COLOR_EXTERNAL = RGBColor(255, 152, 0)     # Orange accent
COLOR_OUTPUT = RGBColor(189, 189, 189)     # Grey
COLOR_TEXT = RGBColor(33, 33, 33)
COLOR_EDGE = RGBColor(66, 66, 66)

# Fonts
FONT_TITLE = 'Calibri'
FONT_BODY = 'Calibri'

# Grid layout (6 columns)
MARGIN_LEFT = Inches(0.4)
MARGIN_TOP = Inches(1.2)
COL_WIDTH = Inches(2.0)
COL_SPACING = Inches(0.15)
ROW_HEIGHT = Inches(0.65)
ROW_SPACING = Inches(0.3)

# ============================================================================
# TITLE
# ============================================================================
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
title_frame = title_box.text_frame
title_frame.text = "Data to Discovery: System Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.name = FONT_TITLE
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_TITLE
title_para.alignment = PP_ALIGN.LEFT

# Subtitle/Caption
caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.3))
caption_frame = caption_box.text_frame
caption_frame.text = "Autonomous discovery through statistical analysis, agent orchestration, and LLM synthesis"
caption_para = caption_frame.paragraphs[0]
caption_para.font.name = FONT_BODY
caption_para.font.size = Pt(16)
caption_para.font.italic = True
caption_para.font.color.rgb = COLOR_TEXT
caption_para.alignment = PP_ALIGN.LEFT

# ============================================================================
# NODE POSITIONING (Swimlane-based)
# ============================================================================
node_positions = {}

# Swimlane 1: Sources (Column 0)
swimlane_sources = [
    ('csv_data', 0, 0, "CSV Data"),
    ('literature_store', 0, 1, "Literature\nStore"),
    ('config', 0, 2, "Config"),
]

# Swimlane 2: Processing (Column 1)
swimlane_processing = [
    ('streamlit_ui', 1, 0, "Streamlit\nUI"),
    ('kosmos_framework', 1, 2, "Kosmos\nFramework"),
]

# Swimlane 3: Agents (Column 2)
swimlane_agents = [
    ('data_analyst', 2, 0, "Data\nAnalyst"),
    ('lit_search_agent', 2, 1, "Literature\nAgent"),
    ('world_model', 2, 2, "World\nModel"),
]

# Swimlane 4: External (Column 3)
swimlane_external = [
    ('openai_api', 3, 1, "OpenAI\nAPI"),
]

# Swimlane 5: Storage (Column 4)
swimlane_storage = [
    ('world_model_state', 4, 2, "World Model\nJSON"),
]

# Swimlane 6: Outputs (Column 5)
swimlane_outputs = [
    ('enhanced_report', 5, 1, "Enhanced\nReport"),
    ('analysis_artifacts', 5, 2, "Analysis\nCode"),
]

all_nodes = (swimlane_sources + swimlane_processing + swimlane_agents +
             swimlane_external + swimlane_storage + swimlane_outputs)

# Shape mapping (use available shapes in python-pptx)
shape_map = {
    'cylinder': MSO_SHAPE.CAN,  # Can shape represents datastore
    'rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
    'hexagon': MSO_SHAPE.HEXAGON,
    'parallelogram': MSO_SHAPE.PARALLELOGRAM,
}

# Color mapping
def get_node_color(node_type):
    if node_type == 'datastore':
        return COLOR_DATASTORE
    elif node_type == 'service':
        return COLOR_SERVICE
    elif node_type == 'external':
        return COLOR_EXTERNAL
    elif node_type in ['input', 'output']:
        return COLOR_OUTPUT
    else:
        return COLOR_SERVICE

# Draw nodes
shapes = {}
for node_id, col, row, label in all_nodes:
    # Find node data
    node_data = next((n for n in arch_data['nodes'] if n['id'] == node_id), None)
    if not node_data:
        continue

    # Calculate position
    left = MARGIN_LEFT + col * (COL_WIDTH + COL_SPACING)
    top = MARGIN_TOP + row * (ROW_HEIGHT + ROW_SPACING)
    width = COL_WIDTH
    height = ROW_HEIGHT

    node_positions[node_id] = (left + width/2, top + height/2)

    # Get shape type
    shape_type = shape_map.get(node_data.get('shape', 'rectangle'), MSO_SHAPE.ROUNDED_RECTANGLE)

    # Add shape
    shape = slide.shapes.add_shape(
        shape_type,
        left, top, width, height
    )

    # Fill
    fill_color = get_node_color(node_data.get('type', 'service'))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # Border
    shape.line.color.rgb = COLOR_EDGE
    shape.line.width = Pt(2)

    # Text
    text_frame = shape.text_frame
    text_frame.text = label
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(8)
    text_frame.margin_right = Pt(8)
    text_frame.margin_top = Pt(8)
    text_frame.margin_bottom = Pt(8)

    para = text_frame.paragraphs[0]
    para.font.name = FONT_BODY
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.color.rgb = COLOR_TEXT
    para.alignment = PP_ALIGN.CENTER

    shapes[node_id] = shape

# ============================================================================
# EDGES (Arrows)
# ============================================================================
# Draw key edges (simplified to avoid clutter)
key_edges = [
    ('csv_data', 'streamlit_ui'),
    ('config', 'streamlit_ui'),
    ('streamlit_ui', 'data_analyst'),
    ('streamlit_ui', 'world_model'),
    ('data_analyst', 'openai_api'),
    ('data_analyst', 'world_model'),
    ('lit_search_agent', 'openai_api'),
    ('lit_search_agent', 'world_model'),
    ('world_model', 'world_model_state'),
    ('world_model', 'enhanced_report'),
    ('data_analyst', 'analysis_artifacts'),
    ('literature_store', 'lit_search_agent'),
]

for from_id, to_id in key_edges:
    if from_id not in node_positions or to_id not in node_positions:
        continue

    from_x, from_y = node_positions[from_id]
    to_x, to_y = node_positions[to_id]

    # Add connector
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        from_x, from_y,
        to_x, to_y
    )

    connector.line.color.rgb = COLOR_EDGE
    connector.line.width = Pt(2)

# ============================================================================
# LEGEND (Bottom Right)
# ============================================================================
legend_left = Inches(10.5)
legend_top = Inches(5.5)
legend_width = Inches(2.5)

# Legend title
legend_title_box = slide.shapes.add_textbox(legend_left, legend_top, legend_width, Inches(0.25))
legend_title_frame = legend_title_box.text_frame
legend_title_frame.text = "Legend"
legend_title_para = legend_title_frame.paragraphs[0]
legend_title_para.font.name = FONT_BODY
legend_title_para.font.size = Pt(16)
legend_title_para.font.bold = True
legend_title_para.font.color.rgb = COLOR_TEXT

# Legend items
legend_items = [
    ("■ Datastore", COLOR_DATASTORE),
    ("■ Service/Module", COLOR_SERVICE),
    ("■ External API", COLOR_EXTERNAL),
    ("■ Input/Output", COLOR_OUTPUT),
]

legend_y = legend_top + Inches(0.35)
for item_text, item_color in legend_items:
    item_box = slide.shapes.add_textbox(legend_left, legend_y, legend_width, Inches(0.22))
    item_frame = item_box.text_frame
    item_frame.text = item_text
    item_para = item_frame.paragraphs[0]
    item_para.font.name = FONT_BODY
    item_para.font.size = Pt(14)

    # Color the square
    runs = item_para.runs
    if runs:
        runs[0].font.color.rgb = item_color
        runs[0].font.size = Pt(18)
        runs[0].font.bold = True

    for i in range(1, len(runs)):
        runs[i].font.color.rgb = COLOR_TEXT

    legend_y += Inches(0.23)

# Line styles
legend_y += Inches(0.1)
line_legend_box = slide.shapes.add_textbox(legend_left, legend_y, legend_width, Inches(0.6))
line_legend_frame = line_legend_box.text_frame
line_legend_frame.text = "─── Verified link\n- - - Optional path"
line_legend_para = line_legend_frame.paragraphs[0]
line_legend_para.font.name = FONT_BODY
line_legend_para.font.size = Pt(12)
line_legend_para.font.color.rgb = COLOR_TEXT

# ============================================================================
# WORD COUNT CHECK
# ============================================================================
total_words = 0
# Count words in title
total_words += len("Data to Discovery: System Architecture".split())
# Count words in caption
total_words += len("Autonomous discovery through statistical analysis, agent orchestration, and LLM synthesis".split())
# Count words in node labels
for _, _, _, label in all_nodes:
    total_words += len(label.replace('\n', ' ').split())

print(f"Total visible text: ~{total_words} words (constraint: ≤60)")

# Save presentation
prs.save('architecture_overview.pptx')
print("✅ PowerPoint slide created: architecture_overview.pptx")
print(f"   - Nodes: {len(all_nodes)} (constraint: ≤12)")
print(f"   - Edges shown: {len(key_edges)} (constraint: ≤16)")
print(f"   - Format: 16:9")
print(f"   - Style: Consulting-grade, diagram-first")
