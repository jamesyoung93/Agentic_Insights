"""
Generate Consulting-Style Architecture Slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_architecture_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Colors
    PRIMARY = RGBColor(0, 32, 96)
    ACCENT = RGBColor(0, 123, 255)
    GRAY_DARK = RGBColor(64, 64, 64)
    GRAY_MED = RGBColor(128, 128, 128)
    GRAY_LIGHT = RGBColor(217, 217, 217)
    GREEN = RGBColor(0, 176, 80)

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title.text_frame.text = "Kosmos AI Scientist: Autonomous Data-Driven Discovery Platform"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12), Inches(0.4))
    sub.text_frame.text = "Transforms customer data into statistically-validated business insights through automated discovery cycles"
    sub.text_frame.paragraphs[0].font.size = Pt(12)
    sub.text_frame.paragraphs[0].font.color.rgb = GRAY_DARK

    # Positions
    col1_x, col2_x, col3_x, col4_x, col5_x = 0.5, 2.5, 4.8, 8.5, 10.8
    base_y = 1.5

    def box(x, y, w, h, txt, shp='RECT', fill=GRAY_LIGHT, txt_color=GRAY_DARK, border=GRAY_DARK):
        if shp == 'CYL':
            s = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DATA, Inches(x), Inches(y), Inches(w), Inches(h))
        elif shp == 'HEX':
            s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(w), Inches(h))
        elif shp == 'PARA':
            s = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(x), Inches(y), Inches(w), Inches(h))
        else:
            s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
        s.text_frame.text = txt
        s.text_frame.word_wrap = True
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in s.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9)
            p.font.color.rgb = txt_color
            p.font.bold = True
        return s

    def arrow(x1, y1, x2, y2, lbl="", dotted=False):
        conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = GRAY_MED
        conn.line.width = Pt(2)
        if dotted:
            conn.line.dash_style = 2
        conn.line.end_arrow_type = 2
        if lbl:
            mid_x, mid_y = (x1 + x2) / 2, (y1 + y2) / 2 - 0.3
            l = slide.shapes.add_textbox(Inches(mid_x - 0.5), Inches(mid_y), Inches(1), Inches(0.3))
            l.text_frame.text = lbl
            l.text_frame.paragraphs[0].font.size = Pt(7)
            l.text_frame.paragraphs[0].font.color.rgb = ACCENT
            l.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return conn

    def group(x, y, w, h, title):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(245, 245, 245)
        b.line.color.rgb = GRAY_LIGHT
        b.line.width = Pt(1)
        t = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.25))
        t.text_frame.text = title
        t.text_frame.paragraphs[0].font.size = Pt(8)
        t.text_frame.paragraphs[0].font.bold = True
        t.text_frame.paragraphs[0].font.color.rgb = PRIMARY
        return b

    # Groups
    group(0.4, base_y - 0.2, 1.9, 4.8, "DATA SOURCES")
    group(2.4, base_y - 0.2, 2.0, 4.8, "INGESTION")
    group(4.6, base_y - 0.2, 3.6, 4.8, "CORE PROCESSING")
    group(8.3, base_y - 0.2, 2.2, 4.8, "STORAGE")
    group(10.6, base_y - 0.2, 2.2, 4.8, "OUTPUTS")

    # Data Sources
    box(col1_x, base_y + 0.3, 1.7, 0.7, "CSV Files\ncustomers.csv\ncompetitor_data.csv", 'CYL', RGBColor(230, 240, 255))
    box(col1_x, base_y + 1.5, 1.7, 0.7, "Literature KB\nknowledge/\nliterature/", 'CYL', RGBColor(230, 240, 255))
    box(col1_x, base_y + 2.7, 1.7, 0.7, "OpenAI API\nGPT-3.5/4\n(Optional)", 'HEX', RGBColor(255, 245, 230))

    # Ingestion
    box(col2_x, base_y + 0.3, 1.7, 1.2, "Data Ingestion\n& Preprocessing\n\nMerge + Clean\nImpute + Engineer", fill=RGBColor(220, 230, 240))

    # Core Processing
    box(col3_x, base_y + 0.1, 1.5, 0.9, "Discovery Cycle\nEngine\n\nOrchestrates\nQ->A->S", fill=ACCENT, txt_color=RGBColor(255, 255, 255), border=ACCENT)
    box(col3_x + 2.0, base_y + 0.1, 1.5, 0.9, "Statistical\nAnalysis\n\nscipy + numpy", fill=RGBColor(220, 230, 240))
    box(col3_x, base_y + 1.3, 1.5, 0.8, "Literature\nSearch Agent\n\nLLM-powered", fill=RGBColor(220, 230, 240))
    box(col3_x + 2.0, base_y + 1.3, 1.5, 0.8, "Discovery\nSynthesis\n\nLLM insights", fill=RGBColor(220, 230, 240))
    box(col3_x + 1.0, base_y + 2.4, 1.5, 0.8, "Report\nGenerator\n\nExtract stats", fill=RGBColor(220, 230, 240))

    # Storage
    box(col4_x, base_y + 0.5, 1.8, 1.5, "World Model\n(Knowledge State)\n\nDiscoveries\nTrajectories\nHypotheses", 'CYL', RGBColor(255, 250, 220))

    # Outputs
    box(col5_x, base_y + 0.1, 1.8, 0.8, "Streamlit\nWeb UI\n\nInteractive", 'PARA', GREEN, RGBColor(255, 255, 255), GREEN)
    box(col5_x, base_y + 1.2, 1.8, 0.8, "Discovery\nReports\n\n.txt files", 'PARA', GREEN, RGBColor(255, 255, 255), GREEN)
    box(col5_x, base_y + 2.3, 1.8, 0.8, "World Model\nJSON\n\nPersistent", 'PARA', GREEN, RGBColor(255, 255, 255), GREEN)

    # Arrows
    arrow(col1_x + 1.7, base_y + 0.65, col2_x, base_y + 0.9, "CSV")
    arrow(col2_x + 1.7, base_y + 0.9, col3_x, base_y + 0.55, "DataFrame")
    arrow(col3_x + 1.5, base_y + 0.35, col3_x + 2.0, base_y + 0.35, "Q")
    arrow(col3_x + 2.0, base_y + 0.75, col3_x + 1.5, base_y + 0.75, "Stats")
    arrow(col3_x + 0.75, base_y + 1.0, col3_x + 0.75, base_y + 1.3, "Query")
    arrow(col3_x + 1.2, base_y + 1.3, col3_x + 1.2, base_y + 1.0, "Results")
    arrow(col3_x + 1.5, base_y + 1.5, col3_x + 2.0, base_y + 1.5, "")
    arrow(col1_x + 1.7, base_y + 3.1, col3_x + 0.5, base_y + 1.0, "LLM", True)
    arrow(col1_x + 1.7, base_y + 1.85, col3_x, base_y + 1.7, "Texts")
    arrow(col3_x + 1.5, base_y + 0.55, col4_x, base_y + 1.25, "Disc")
    arrow(col3_x + 3.5, base_y + 1.7, col4_x, base_y + 1.25, "")
    arrow(col4_x, base_y + 2.0, col3_x + 2.5, base_y + 2.8, "Data")
    arrow(col3_x + 2.5, base_y + 2.8, col5_x, base_y + 1.6, "Format")
    arrow(col4_x + 1.8, base_y + 1.75, col5_x, base_y + 2.7, "JSON")
    arrow(col4_x + 1.8, base_y + 0.75, col5_x, base_y + 0.5, "State")
    arrow(col5_x, base_y + 0.2, col3_x + 0.75, base_y - 0.1, "Config", True)

    # Legend
    lx, ly = 0.5, 6.3
    lt = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(2), Inches(0.25))
    lt.text_frame.text = "LEGEND"
    lt.text_frame.paragraphs[0].font.size = Pt(9)
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.color.rgb = PRIMARY

    items = [
        ("Cylinder = Data Store", 'CYL', RGBColor(230, 240, 255)),
        ("Rectangle = Processing", 'RECT', RGBColor(220, 230, 240)),
        ("Hexagon = External API", 'HEX', RGBColor(255, 245, 230)),
        ("Parallelogram = Output", 'PARA', GREEN)
    ]

    for i, (txt, shp, col) in enumerate(items):
        y = ly + 0.3 + (i * 0.35)
        if shp == 'CYL':
            s = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DATA, Inches(lx), Inches(y), Inches(0.25), Inches(0.2))
        elif shp == 'HEX':
            s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(lx), Inches(y), Inches(0.25), Inches(0.2))
        elif shp == 'PARA':
            s = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(lx), Inches(y), Inches(0.25), Inches(0.2))
        else:
            s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(y), Inches(0.25), Inches(0.2))
        s.fill.solid()
        s.fill.fore_color.rgb = col
        s.line.color.rgb = GRAY_DARK
        s.line.width = Pt(1)
        l = slide.shapes.add_textbox(Inches(lx + 0.35), Inches(y), Inches(1.5), Inches(0.2))
        l.text_frame.text = txt
        l.text_frame.paragraphs[0].font.size = Pt(7)
        l.text_frame.paragraphs[0].font.color.rgb = GRAY_DARK
        l.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Runtime info
    rt = slide.shapes.add_textbox(Inches(lx + 2.5), Inches(ly), Inches(3.5), Inches(1.5))
    rt.text_frame.text = "RUNTIME TRIGGERS\n- User initiates via Streamlit UI\n- Cycles run sequentially (1-20)\n- Per-cycle: 3 questions\n- LLM calls: Optional"
    for p in rt.text_frame.paragraphs:
        p.font.size = Pt(7)
        p.font.color.rgb = GRAY_DARK
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.size = Pt(9)
    rt.text_frame.paragraphs[0].font.color.rgb = PRIMARY

    # Deployment
    dep = slide.shapes.add_textbox(Inches(lx + 6.5), Inches(ly), Inches(3), Inches(1.5))
    dep.text_frame.text = "DEPLOYMENT\n- Local Python (streamlit run)\n- Deps: pandas, scipy, streamlit\n- Optional: OpenAI API key\n- Scale: Single-user, local"
    for p in dep.text_frame.paragraphs:
        p.font.size = Pt(7)
        p.font.color.rgb = GRAY_DARK
    dep.text_frame.paragraphs[0].font.bold = True
    dep.text_frame.paragraphs[0].font.size = Pt(9)
    dep.text_frame.paragraphs[0].font.color.rgb = PRIMARY

    # Save
    output = "architecture_overview.pptx"
    prs.save(output)
    print(f"Created: {output}")
    return output

if __name__ == "__main__":
    create_architecture_slide()
